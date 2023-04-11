import sys
import re
from pptx import Presentation
from pptx.util import Inches
from markdown import markdown
from bs4 import BeautifulSoup

def create_slide(pres, title, content, tables):
    slide_layout = pres.slide_layouts[1]
    slide = pres.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    title_shape.text = title

    soup = BeautifulSoup(content, 'html.parser')
    content = soup.get_text('\n')

    paragraphs = content.split('\n')

    text_frame = content_shape.text_frame
    text_frame.clear()

    for paragraph in paragraphs:
        stripped_paragraph = paragraph.strip()
        if not stripped_paragraph:
            continue
        p = text_frame.add_paragraph()
        run = p.add_run()
        run.text = stripped_paragraph

    for table_data in tables:
        add_table_to_slide(slide, table_data)

    return slide


def add_table_to_slide(slide, table_data):
    rows, cols = len(table_data), len(table_data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(4))
    table = table_shape.table

    for row_idx, row in enumerate(table_data):
        for col_idx, value in enumerate(row):
            table.cell(row_idx, col_idx).text = str(value)

def parse_tables(html_content):
    soup = BeautifulSoup(html_content, 'html.parser')
    tables = []

    for table_tag in soup.find_all('table'):
        table = []
        for row_tag in table_tag.find_all('tr'):
            row = []
            for cell_tag in row_tag.find_all(['td', 'th']):
                row.append(cell_tag.get_text())
            table.append(row)
        tables.append(table)
        table_tag.decompose()

    return tables, str(soup)

def main(input_file, output_file):
    with open(input_file, 'r') as file:
        md_content = file.read()

    html_content = markdown(md_content, extensions=['tables'])
    slides = re.split(r'<h1>(.*?)</h1>', html_content)
    slides.pop(0)

    pres = Presentation()

    for i in range(0, len(slides), 2):
        content = slides[i + 1]
        tables, content = parse_tables(content)
        slide = create_slide(pres, slides[i], content, tables)

    pres.save(output_file)

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python markdown_to_keynote.py input.md output.pptx")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2]
    main(input_file, output_file)
